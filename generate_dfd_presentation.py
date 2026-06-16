import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_dfd_presentation():
    prs = Presentation()
    
    # Set slide width and height to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Helper function to create image slides with title and text
    def add_image_text_slide(title_text, image_path, explanation_points):
        # We will use a blank layout and add things manually for better control
        layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(layout)
        
        # Add Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.bold = True
        
        # Add Image (Left side)
        img_width = Inches(7.5)
        if os.path.exists(image_path):
            left = Inches(0.5)
            top = Inches(1.2)
            slide.shapes.add_picture(image_path, left, top, width=img_width)
        else:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(7.5), Inches(5))
            txBox.text_frame.text = f"Image not found at: {image_path}"

        # Add Text Explanation (Right side)
        text_box = slide.shapes.add_textbox(Inches(8.2), Inches(1.2), Inches(4.8), Inches(5.5))
        tf_exp = text_box.text_frame
        tf_exp.word_wrap = True
        
        for i, point in enumerate(explanation_points):
            if i == 0:
                p_exp = tf_exp.paragraphs[0]
            else:
                p_exp = tf_exp.add_paragraph()
            p_exp.text = point
            p_exp.font.size = Pt(20)
            p_exp.level = 0
            p_exp.space_after = Pt(14)

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Data Flow Diagrams (DFD)"
    subtitle.text = "TravelEase Project\nSystem Architecture & Data Movement"

    # Add DFD Slides
    base_dir = "synopsis_diagrams"
    
    # DFD Level 0
    add_image_text_slide(
        "Context Diagram (DFD Level 0)",
        os.path.join(base_dir, "dfd_level_0.png"),
        [
            "Provides a high-level view of the entire TravelEase system.",
            "Shows interactions between the central system and external entities (User, Admin, AI API).",
            "Data flows in as user requests (searches, bookings, logins).",
            "Data flows out as responses (booking confirmations, chatbot answers, dashboard data)."
        ]
    )
    
    # DFD Level 1
    add_image_text_slide(
        "System Processes (DFD Level 1)",
        os.path.join(base_dir, "dfd_level_1.png"),
        [
            "Breaks down the main system into major functional subsystems.",
            "Key processes include Authentication, Package Booking, Transport Booking, and Admin Management.",
            "Illustrates how data stores (Databases) are accessed by different modules.",
            "Shows the independent flow of transport logic vs package logic."
        ]
    )
    
    # DFD Level 2
    add_image_text_slide(
        "Detailed Sub-Processes (DFD Level 2)",
        os.path.join(base_dir, "dfd_level_2.png"),
        [
            "Dives deep into specific complex processes like Transport Booking or Checkout.",
            "Shows granular steps: e.g., Seat selection -> Payment Validation -> DB Update -> Receipt.",
            "Clarifies validation checks and data transformation at the micro-level.",
            "Essential for understanding the actual implementation logic of core features."
        ]
    )

    # Save the presentation
    output_filename = "TravelEase_DFD_Slides_v2.pptx"
    prs.save(output_filename)
    print(f"Successfully generated {output_filename}")

if __name__ == "__main__":
    create_dfd_presentation()
