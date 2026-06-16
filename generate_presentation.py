import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Optional: Set slide width and height to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "TravelEase\nWeb-Based Holiday and Transport Booking System"
    subtitle.text = "MCA Major Project\n\nSubmitted by: [Your Name]\nRoll No: [Your Roll Number]"

    # Helper function to create content slides
    def add_bullet_slide(title_text, bullet_points):
        layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(layout)
        title_shape = slide.shapes.title
        title_shape.text = title_text
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(24)
            p.level = 0
            
    # 2. Introduction
    add_bullet_slide(
        "Introduction",
        [
            "The travel sector relies heavily on digital booking platforms.",
            "Users currently face a fragmented experience (separate apps for flights, trains, packages).",
            "TravelEase is a unified solution integrating package discovery, transport booking, and user management.",
            "Designed as a full-stack academic project using Python and Flask."
        ]
    )

    # 3. Problem Statement
    add_bullet_slide(
        "Problem Statement",
        [
            "Fragmented Booking Experience: No single platform for comprehensive holiday planning.",
            "Lack of Unified History: Difficult to track bookings across multiple services.",
            "Weak Admin Controls: Small systems lack proper UI-based management.",
            "Insufficient Decision Support: Missing features like chat support, reviews, or wishlists."
        ]
    )

    # 4. Proposed Solution
    add_bullet_slide(
        "Proposed Solution: TravelEase",
        [
            "Centralized Platform: Holiday packages, flights, trains, buses, and cabs in one place.",
            "Unified User Dashboard: All booking types tracked centrally.",
            "Role-Based Admin Panel: Complete CRUD operations via a clean UI.",
            "Decision-Support Tools: AI/DB-aware chatbot, package reviews, and wishlist.",
            "Validated Checkout: Proper payment validation before booking confirmation."
        ]
    )

    # 5. Objectives
    add_bullet_slide(
        "Objectives of the Project",
        [
            "Develop an end-to-end holiday booking and multi-mode transport system.",
            "Implement secure, role-based user authentication (User vs Admin).",
            "Maintain a normalized relational database schema (MySQL via SQLAlchemy).",
            "Provide engagement features like daily spin-wheel discounts.",
            "Ensure maintainability and Docker-ready deployment."
        ]
    )

    # 6. System Architecture
    add_bullet_slide(
        "System Architecture",
        [
            "3-Tier Client-Server Architecture:",
            "  â€¢ Presentation Tier: HTML (Jinja2), CSS, JavaScript.",
            "  â€¢ Application Tier: Flask (Python), handling routes, logic, and auth.",
            "  â€¢ Data Tier: MySQL / SQLite (fallback) managed by SQLAlchemy ORM.",
            "Session-based state management for discounts and search."
        ]
    )

    # 7. Technology Stack
    add_bullet_slide(
        "Technology Stack",
        [
            "Backend Framework: Python with Flask.",
            "Database: MySQL 8.0 (Primary) / SQLite (Fallback).",
            "ORM & DB Tools: Flask-SQLAlchemy, PyMySQL.",
            "Frontend: HTML5, CSS3, JavaScript, Jinja2 Templates.",
            "Deployment: Docker (containerization).",
            "Security: Werkzeug for password hashing (scrypt)."
        ]
    )

    # 8. Key Features: User Panel
    add_bullet_slide(
        "Key Features: User Panel",
        [
            "User Registration, Login, and secure session management.",
            "Holiday Package Browsing (filtering by price/keyword).",
            "Transport Booking (Flights, Trains, Buses, Cabs) with class-based pricing.",
            "Unified Dashboard to view all historical bookings.",
            "Wishlist and Review system (1 review per package)."
        ]
    )

    # 9. Key Features: Admin Panel
    add_bullet_slide(
        "Key Features: Admin Panel",
        [
            "Accessible only to users with 'is_admin = True'.",
            "Package Management: Add, edit, and delete holiday packages.",
            "Transport Management: Manage available flights, trains, buses, cabs.",
            "Booking Oversight: View all user bookings across the platform.",
            "Customer Support: Acknowledge and track user inquiries."
        ]
    )

    # 10. Engagement & Unique Features
    add_bullet_slide(
        "Engagement & Unique Features",
        [
            "Daily Spin Wheel: Generates session-based discount codes for checkout.",
            "Fallback Database Mechanism: Auto-switches to SQLite if MySQL is down.",
            "Payment Validation: Validates card formats (number, expiry, CVV) or UPI.",
            "Interactive UI: Dynamic components for a modern user experience."
        ]
    )

    # 11. Chatbot Integration
    add_bullet_slide(
        "3-Layer Chatbot Architecture",
        [
            "Layer 1 (Database-Aware): Matches intents (budget, transport) and queries live DB.",
            "Layer 2 (AI-Powered): Uses OpenAI API to generate contextual responses based on packages.",
            "Layer 3 (Structured Fallback): Keyword-based static responses if DB/AI fails.",
            "Result: Ensures accurate, reliable, and helpful 24/7 user support."
        ]
    )

    # 12. Database Design
    add_bullet_slide(
        "Database Design (Schema)",
        [
            "Built on Third Normal Form (3NF) to minimize redundancy.",
            "Core Entities: Users, Packages, Bookings (Generic & Transport-specific).",
            "Engagement Entities: Reviews, Wishlist, Contact Inquiries.",
            "Key Relationships: Bookings link to both Users and Packages/Transport via Foreign Keys."
        ]
    )

    # 13. Deployment Strategy
    add_bullet_slide(
        "Deployment Strategy",
        [
            "Docker Containerization: Simplifies setup using Dockerfile and docker-compose.",
            "Automated Database Seeding: Scripts automatically populate tables on startup.",
            "Environment Agnostic: Uses .env files for configuration (DB URLs, API keys).",
            "Production Ready: Modular codebase easily extensible to a WSGI server."
        ]
    )

    # 14. Future Scope
    add_bullet_slide(
        "Future Enhancements",
        [
            "Integration with real live Payment Gateways (e.g., Stripe, Razorpay).",
            "Real-time API sync for live flight and train availability.",
            "Email or SMS notifications for booking confirmations.",
            "Multi-language support for broader accessibility.",
            "Dedicated Mobile Application interface."
        ]
    )

    # 15. Conclusion & Q&A
    add_bullet_slide(
        "Conclusion",
        [
            "TravelEase successfully demonstrates a unified travel booking system.",
            "It bridges the gap between fragmented transport and package bookings.",
            "Provides a complete academic-grade implementation of web dev principles.",
            "Thank You! Questions?"
        ]
    )

    # Save the presentation
    output_filename = "TravelEase_Presentation.pptx"
    prs.save(output_filename)
    print(f"Successfully generated {output_filename}")

if __name__ == "__main__":
    create_presentation()
