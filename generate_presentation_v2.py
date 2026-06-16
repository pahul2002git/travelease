import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_beautiful_presentation():
    prs = Presentation()
    
    # Set to 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Set background to dark blue (glassmorphism vibe)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate 900
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        for p in subtitle.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(148, 163, 184) # Slate 400

    def add_content_slide_with_image(title_text, bullet_points, image_path=None, layout_type="right_image"):
        # We use a blank layout and construct elements manually for better control
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Dark Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)
        
        # Add Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.color.rgb = RGBColor(56, 189, 248) # Light Blue 400
        p.font.size = Pt(36)
        p.font.bold = True

        # Text Box positioning based on layout
        if image_path and os.path.exists(image_path):
            if layout_type == "right_image":
                text_width = Inches(6.0)
                text_left = Inches(0.5)
                img_left = Inches(7.0)
                img_width = Inches(5.8)
            else:
                text_width = Inches(12.0)
                text_left = Inches(0.5)
                img_left = Inches(2.0)
                img_width = Inches(9.333)
                
            # Add image
            try:
                # Calculate image height proportionally if possible, or just let python-pptx scale it
                if layout_type == "right_image":
                    slide.shapes.add_picture(image_path, img_left, Inches(1.8), width=img_width)
                else:
                    # Center image below text
                    slide.shapes.add_picture(image_path, img_left, Inches(2.5), width=img_width)
            except Exception as e:
                print(f"Warning: Could not add image {image_path}: {e}")
        else:
            text_width = Inches(12.0)
            text_left = Inches(0.5)

        # Add Text Content
        body_box = slide.shapes.add_textbox(text_left, Inches(1.8), text_width, Inches(5.0))
        btf = body_box.text_frame
        btf.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = btf.paragraphs[0]
            else:
                p = btf.add_paragraph()
            p.text = "â€¢ " + point
            p.font.color.rgb = RGBColor(241, 245, 249) # Slate 100
            p.font.size = Pt(22)
            p.space_after = Pt(14)


    # 1. Title Slide
    add_title_slide(
        "TravelEase",
        "Web-Based Holiday and Transport Booking System\n\nMCA Major Project\nSubmitted by: [Your Name]\nRoll No: [Your Roll Number]"
    )

    # 2. Introduction
    add_content_slide_with_image(
        "Project Introduction",
        [
            "The travel sector relies heavily on digital booking platforms.",
            "Users currently face a fragmented experience across separate apps.",
            "TravelEase is a unified solution integrating package discovery, transport booking, and user management.",
            "Built as a full-stack web application using modern design principles."
        ],
        "screenshots/home.png", 
        "right_image"
    )

    # 3. System Architecture
    add_content_slide_with_image(
        "System Architecture",
        [
            "3-Tier Client-Server Architecture:",
            "Presentation Tier: HTML (Jinja2), CSS, JavaScript.",
            "Application Tier: Flask (Python), handling routes, logic, and auth.",
            "Data Tier: MySQL managed by SQLAlchemy ORM."
        ],
        "final_report_assets/Fig_01_System_Architecture.png",
        "bottom_image"
    )

    # 4. Database Design (ER)
    add_content_slide_with_image(
        "Database Design",
        [
            "Built on Third Normal Form (3NF) to minimize redundancy.",
            "Central Users, Packages, and Bookings tables.",
            "Foreign Key relationships maintain strict data integrity."
        ],
        "final_report_assets/Fig_05_ER_Diagram.png",
        "right_image"
    )

    # 5. User Features: Home & Registration
    add_content_slide_with_image(
        "User Registration & Login",
        [
            "Secure user registration with password hashing (scrypt).",
            "Session-based state management.",
            "Clean, glassmorphism-inspired UI for accessibility."
        ],
        "screenshots/login.png",
        "right_image"
    )
    
    # 6. User Features: Packages
    add_content_slide_with_image(
        "Package Browsing & Booking",
        [
            "Dynamic filtering by price and keywords.",
            "Detailed package views including itineraries and galleries.",
            "Integrated Checkout with secure payment validation."
        ],
        "screenshots/package_detail.png",
        "right_image"
    )

    # 7. User Features: Dashboard
    add_content_slide_with_image(
        "Unified User Dashboard",
        [
            "Centralized hub for all booking types.",
            "View past flights, trains, buses, cabs, and holidays.",
            "Includes wishlist management and package reviews."
        ],
        "screenshots/dashboard.png",
        "right_image"
    )

    # 8. User Features: Transport
    add_content_slide_with_image(
        "Multi-Mode Transport Booking",
        [
            "Dedicated modules for Flights, Trains, Buses, and Cabs.",
            "Class-based pricing (Economy, Business, First Class).",
            "Automatic fare computation dynamically updated."
        ],
        "screenshots/flights.png",
        "right_image"
    )
    
    # 9. Admin Operations
    add_content_slide_with_image(
        "Role-Based Admin Panel",
        [
            "Accessible only to users with 'is_admin = True'.",
            "Perform full CRUD operations on packages and transport.",
            "Oversee all user bookings and customer inquiries."
        ],
        "final_report_assets/Fig_10_Flow_Admin_Operations.png",
        "right_image"
    )

    # 10. AI Chatbot Integration
    add_content_slide_with_image(
        "3-Layer Chatbot Architecture",
        [
            "Layer 1: Database-Aware keyword matching for live pricing.",
            "Layer 2: AI-Powered fallback using OpenAI API.",
            "Layer 3: Structured Fallback for offline resilience."
        ],
        "final_report_assets/Fig_11_Flow_Chatbot_Working.png",
        "right_image"
    )

    # 11. Conclusion
    add_content_slide_with_image(
        "Conclusion & Future Scope",
        [
            "TravelEase successfully demonstrates a unified travel ecosystem.",
            "Bridged the gap between fragmented booking services.",
            "Future Scope: Integration with live Payment Gateways (Stripe).",
            "Future Scope: Real-time API sync for live flight availability.",
            "Thank You! Questions?"
        ],
        None,
        "full_text"
    )

    # Save the presentation
    output_filename = "TravelEase_Presentation_v2.pptx"
    prs.save(output_filename)
    print(f"Successfully generated {output_filename}")

if __name__ == "__main__":
    create_beautiful_presentation()
